import json
import os
import re
import importlib.util
from io import BytesIO, StringIO
from typing import Any, Dict, List

import streamlit as st

from prompts import OUTPUT_SCHEMA, PROMPT_LIBRARY, SYSTEM_PROMPT, USER_PROMPT_TEMPLATE


DEFAULT_MODEL = "gemini-2.5-flash"
KNOWN_OWNERS = {"hr", "it", "manager", "supervisor", "operations", "finance", "team lead", "security"}

OPTIONAL_DEPENDENCIES = {
    "google-genai": "google.genai",
    "pypdf": "pypdf",
    "python-docx": "docx",
    "python-pptx": "pptx",
}


def get_prompt_config(prompt_version: str | None = None) -> Dict[str, str]:
    if prompt_version and prompt_version in PROMPT_LIBRARY:
        return PROMPT_LIBRARY[prompt_version]

    if "v2" in PROMPT_LIBRARY:
        return PROMPT_LIBRARY["v2"]

    if PROMPT_LIBRARY:
        return next(iter(PROMPT_LIBRARY.values()))

    return {
        "system_prompt": SYSTEM_PROMPT,
        "output_schema": OUTPUT_SCHEMA,
        "user_prompt_template": USER_PROMPT_TEMPLATE,
    }


def is_gemini_available() -> bool:
    return bool(get_gemini_api_key())


def get_gemini_api_key() -> str:
    secret_value = st.secrets.get("GEMINI_API_KEY", "")
    if secret_value:
        return str(secret_value)
    return os.getenv("GEMINI_API_KEY", "")


def get_missing_dependencies() -> List[str]:
    missing = []
    for package_name, module_name in OPTIONAL_DEPENDENCIES.items():
        if importlib.util.find_spec(module_name) is None:
            missing.append(package_name)
    return missing


def build_dependency_install_message(missing_dependencies: List[str]) -> str:
    packages = ", ".join(missing_dependencies)
    return (
        "Some required Python packages are missing: "
        f"{packages}. Install the project requirements and rerun the app."
    )


def normalize_sop_text(text: str) -> str:
    if not text:
        return ""

    cleaned = text.replace("\r\n", "\n").replace("\r", "\n")
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    lines = [line.strip(" -\t") for line in cleaned.splitlines()]
    normalized_lines = [line for line in lines if line]
    return "\n".join(normalized_lines).strip()


def normalize_inline_text(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def strip_markdown(text: str) -> str:
    cleaned = normalize_inline_text(text)
    cleaned = re.sub(r"[*_`#>\[\]]", "", cleaned)
    return cleaned.strip(" -:")


def sentence_case(text: str) -> str:
    cleaned = strip_markdown(text)
    if not cleaned:
        return ""
    return cleaned[0].upper() + cleaned[1:]


def looks_meaningful(text: str) -> bool:
    cleaned = normalize_inline_text(text)
    return bool(cleaned and re.search(r"[A-Za-z0-9]", cleaned))


def extract_text_from_upload(uploaded_file) -> str:
    file_name = uploaded_file.name.lower()

    if file_name.endswith(".txt"):
        string_io = StringIO(uploaded_file.getvalue().decode("utf-8", errors="ignore"))
        return normalize_sop_text(string_io.read())

    if file_name.endswith(".pdf"):
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise RuntimeError("PDF processing is unavailable right now.") from exc

        try:
            reader = PdfReader(BytesIO(uploaded_file.getvalue()))
            pages = [page.extract_text() or "" for page in reader.pages]
        except Exception as exc:
            raise ValueError("We couldn't read that PDF. Please upload a valid PDF file.") from exc

        extracted = normalize_sop_text("\n".join(pages))
        if not extracted:
            raise ValueError("The PDF did not contain readable text.")
        return extracted

    if file_name.endswith(".docx"):
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("DOCX processing is unavailable right now.") from exc

        try:
            document = Document(BytesIO(uploaded_file.getvalue()))
        except Exception as exc:
            raise ValueError("We couldn't read that DOCX file. Please upload a valid document.") from exc

        paragraphs = [
            paragraph.text
            for paragraph in document.paragraphs
            if looks_meaningful(paragraph.text)
        ]
        table_cells = []
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if looks_meaningful(cell.text):
                        table_cells.append(cell.text)
        extracted = normalize_sop_text("\n".join(paragraphs + table_cells))
        if not extracted:
            raise ValueError("The DOCX file did not contain readable text.")
        return extracted

    raise ValueError("Unsupported file type. Please upload a TXT, PDF, or DOCX file.")


def extract_response_text(response: Any) -> str:
    text = getattr(response, "text", None)
    if isinstance(text, str) and text.strip():
        return text
    raise ValueError("The AI response was empty.")


def safe_json_loads(raw_text: str) -> Dict[str, Any]:
    try:
        return json.loads(raw_text)
    except json.JSONDecodeError:
        start = raw_text.find("{")
        end = raw_text.rfind("}")
        if start != -1 and end != -1:
            return json.loads(raw_text[start : end + 1])
        raise ValueError("The AI response could not be formatted.")


def dedupe_strings(items: List[str], limit: int | None = None) -> List[str]:
    seen = set()
    cleaned_items = []
    for item in items:
        normalized = re.sub(r"\s+", " ", item).strip()
        key = normalized.lower()
        if not normalized or key in seen:
            continue
        seen.add(key)
        cleaned_items.append(normalized)
        if limit is not None and len(cleaned_items) >= limit:
            break
    return cleaned_items


def clip_words(text: str, max_words: int) -> str:
    words = normalize_inline_text(text).split()
    if len(words) <= max_words:
        return " ".join(words)
    return " ".join(words[:max_words]).rstrip(",.;:") + "..."


def normalize_owner(owner: str) -> str:
    cleaned = sentence_case(owner)
    if not cleaned:
        return "Owner"

    lowered = cleaned.lower()
    for known_owner in KNOWN_OWNERS:
        if known_owner in lowered:
            return sentence_case(known_owner)
    return cleaned


def ensure_leading_verb(title: str, action: str) -> str:
    cleaned_title = sentence_case(title)
    if not cleaned_title:
        cleaned_title = sentence_case(action)

    if not cleaned_title:
        return ""

    first_word = cleaned_title.split()[0].lower()
    if first_word.endswith("ing"):
        cleaned_title = f"Complete {cleaned_title.lower()}"
    return cleaned_title


def normalize_decision_trigger(decision_trigger: str) -> str:
    cleaned = sentence_case(decision_trigger)
    if not cleaned:
        return ""
    if "->" in cleaned:
        return cleaned
    if cleaned.lower().startswith("if "):
        return cleaned
    return f"If {cleaned[0].lower() + cleaned[1:]} -> escalate or follow the defined next step"


def remove_summary_overlap(summary: List[str], training_steps: List[Dict[str, str]]) -> List[str]:
    training_signatures = {
        re.sub(r"[^a-z0-9 ]", "", f"{step['owner']} {step['title']} {step['action']}".lower())
        for step in training_steps
    }

    filtered = []
    for bullet in summary:
        signature = re.sub(r"[^a-z0-9 ]", "", bullet.lower())
        if any(signature and signature in item for item in training_signatures):
            continue
        filtered.append(bullet)
    return filtered[:5] or summary[:5]


def ensure_process_improvement_coverage(improvements: List[Dict[str, str]]) -> List[Dict[str, str]]:
    required = {
        "efficiency": {
            "theme": "Efficiency",
            "improvement": "Remove duplicate approvals and batch routine handoffs.",
            "impact": "Reduces delays and shortens cycle time.",
        },
        "automation": {
            "theme": "Automation",
            "improvement": "Automate status alerts, reminders, or system checks.",
            "impact": "Cuts manual follow-up and reduces missed steps.",
        },
        "clarity": {
            "theme": "Clarity",
            "improvement": "Define one owner, trigger, and success check for each step.",
            "impact": "Reduces confusion during execution.",
        },
    }

    present = set()
    normalized_items = []
    for item in improvements:
        theme_key = item["theme"].strip().lower()
        for required_key in required:
            if required_key in theme_key:
                present.add(required_key)
                normalized_items.append(
                    {
                        "theme": sentence_case(required_key),
                        "improvement": sentence_case(item["improvement"]),
                        "impact": sentence_case(item["impact"]),
                    }
                )
                break

    for required_key, fallback in required.items():
        if required_key not in present:
            normalized_items.append(fallback)

    return normalized_items[:3]


def ensure_slide_slots(slides: List[Dict[str, Any]], package: Dict[str, Any]) -> List[Dict[str, Any]]:
    indexed = {slide["slide_number"]: slide for slide in slides if 1 <= slide["slide_number"] <= 8}
    fallback = build_fallback_slide_deck(package.get("sop_name", "SOP"), package)["slides"]
    for fallback_slide in fallback:
        indexed.setdefault(fallback_slide["slide_number"], fallback_slide)
    return [indexed[number] for number in range(1, 9)]


def sanitize_string_list(value: Any, limit: int | None = None) -> List[str]:
    if not isinstance(value, list):
        return []
    return dedupe_strings([str(item).strip() for item in value if str(item).strip()], limit=limit)


def sanitize_summary(summary: Any) -> List[str]:
    cleaned = []
    for item in sanitize_string_list(summary, limit=5):
        text = clip_words(sentence_case(item), 14)
        if text:
            cleaned.append(text)
    return cleaned[:5]


def sanitize_training_steps(training_steps: Any) -> List[Dict[str, str]]:
    if not isinstance(training_steps, list):
        return []

    clean_steps = []
    for item in training_steps:
        if not isinstance(item, dict):
            continue

        title = ensure_leading_verb(str(item.get("title", "")).strip(), str(item.get("action", "")).strip())
        owner = normalize_owner(str(item.get("owner", "")).strip())
        action = sentence_case(str(item.get("action", "")).strip())
        outcome = sentence_case(str(item.get("outcome", "")).strip())
        decision_trigger = normalize_decision_trigger(str(item.get("decision_trigger", "")).strip())

        if not title or not owner or not action or not outcome:
            continue

        clean_steps.append(
            {
                "title": title,
                "owner": owner,
                "action": action,
                "outcome": outcome,
                "decision_trigger": decision_trigger,
            }
        )
        if len(clean_steps) >= 10:
            break

    return clean_steps


def sanitize_quiz(quiz: Any) -> List[Dict[str, str]]:
    if not isinstance(quiz, list):
        return []

    clean_quiz = []
    for item in quiz:
        if not isinstance(item, dict):
            continue

        question = sentence_case(str(item.get("question", "")).strip())
        answer = sentence_case(str(item.get("answer", "")).strip())
        explanation = sentence_case(str(item.get("explanation", "")).strip())

        if question and answer and explanation:
            clean_quiz.append(
                {
                    "question": question,
                    "answer": answer,
                    "explanation": explanation,
                }
            )
        if len(clean_quiz) >= 4:
            break

    return clean_quiz


def sanitize_process_improvements(process_improvements: Any) -> List[Dict[str, str]]:
    if not isinstance(process_improvements, list):
        return []

    clean_improvements = []
    for item in process_improvements:
        if not isinstance(item, dict):
            continue

        theme = sentence_case(str(item.get("theme", "")).strip())
        improvement = sentence_case(str(item.get("improvement", "")).strip())
        impact = sentence_case(str(item.get("impact", "")).strip())

        if theme and improvement and impact:
            clean_improvements.append(
                {
                    "theme": theme,
                    "improvement": improvement,
                    "impact": impact,
                }
            )
        if len(clean_improvements) >= 3:
            break

    return clean_improvements


def sanitize_slide_deck(slides_payload: Any, fallback_title: str) -> Dict[str, Any]:
    if isinstance(slides_payload, list):
        slides_payload = {
            "deck_title": fallback_title,
            "slides": slides_payload,
            "design_notes": {},
        }

    if not isinstance(slides_payload, dict):
        return {
            "deck_title": fallback_title,
            "slides": [],
            "design_notes": {
                "color_theme": "",
                "layout_suggestions": [],
                "font_pairing": "",
            },
        }

    slides = []
    for item in slides_payload.get("slides", []):
        if not isinstance(item, dict):
            continue

        title = sentence_case(str(item.get("title", "")).strip())
        slide_number = item.get("slide_number")
        bullets = [sentence_case(bullet) for bullet in sanitize_string_list(item.get("bullets", []), limit=5)]
        visual_suggestions = [
            sentence_case(visual) for visual in sanitize_string_list(item.get("visual_suggestions", []), limit=3)
        ]

        if title and isinstance(slide_number, int) and bullets:
            slides.append(
                {
                    "slide_number": slide_number,
                    "title": title,
                    "bullets": bullets,
                    "visual_suggestions": visual_suggestions,
                }
            )

    slides = sorted(slides, key=lambda item: item["slide_number"])[:8]

    design_notes = slides_payload.get("design_notes", {})
    if not isinstance(design_notes, dict):
        design_notes = {}

    return {
        "deck_title": str(slides_payload.get("deck_title", "")).strip() or fallback_title,
        "slides": slides,
        "design_notes": {
            "color_theme": sentence_case(str(design_notes.get("color_theme", "")).strip()),
            "layout_suggestions": [
                sentence_case(item)
                for item in sanitize_string_list(design_notes.get("layout_suggestions", []), limit=4)
            ],
            "font_pairing": sentence_case(str(design_notes.get("font_pairing", "")).strip()),
        },
    }


def build_fallback_slide_deck(sop_name: str, package: Dict[str, Any]) -> Dict[str, Any]:
    training_steps = package.get("training_steps", [])
    grouped_steps = [training_steps[0:2], training_steps[2:4], training_steps[4:6], training_steps[6:8]]
    slides = [
        {
            "slide_number": 1,
            "title": sop_name or "SOP Training Deck",
            "bullets": ["Operational training deck", "Executive-ready summary", "Execution-focused content"],
            "visual_suggestions": ["Title divider", "Workflow icon"],
        },
        {
            "slide_number": 2,
            "title": "Overview",
            "bullets": package.get("summary", [])[:5],
            "visual_suggestions": ["Checklist icon", "Section divider"],
        },
    ]

    for index, step_group in enumerate(grouped_steps, start=3):
        bullets = []
        for step in step_group:
            bullets.append(f"{step['owner']}: {step['title']} -> {step['outcome']}")
        if not bullets:
            bullets = ["Complete the next execution checkpoint and confirm owner accountability."]
        slides.append(
            {
                "slide_number": index,
                "title": f"Key Steps {index - 2}",
                "bullets": bullets[:5],
                "visual_suggestions": ["Process flow", "Checklist icon"],
            }
        )

    slides.append(
        {
            "slide_number": 7,
            "title": "Common Mistakes",
            "bullets": package.get("common_mistakes", [])[:5]
            or ["Missed handoffs create avoidable delays and rework."],
            "visual_suggestions": ["Alert icon", "Risk divider"],
        }
    )
    slides.append(
        {
            "slide_number": 8,
            "title": "Improvements",
            "bullets": [
                f"{item['theme']}: {item['improvement']}" for item in package.get("process_improvements", [])
            ][:5]
            or ["Prioritize one process improvement that reduces delay or confusion."],
            "visual_suggestions": ["Automation icon", "Modern closing slide"],
        }
    )

    return {
        "deck_title": f"{sop_name or 'SOP'} Training Deck",
        "slides": slides[:8],
        "design_notes": {
            "color_theme": "Slate, white, and muted teal with minimal accent color.",
            "layout_suggestions": [
                "Use wide margins and one message per slide.",
                "Pair each section with a simple icon and clean divider.",
                "Use left-aligned headings with short bullet stacks.",
            ],
            "font_pairing": "Aptos Display with Aptos or Calibri body text.",
        },
    }


def generate_gemini_training_package(
    sop_text: str,
    model: str = DEFAULT_MODEL,
    prompt_version: str | None = None,
) -> Dict[str, Any]:
    try:
        from google import genai
        from google.genai import types
    except ImportError as exc:
        raise RuntimeError("AI generation is unavailable right now.") from exc

    prompt_config = get_prompt_config(prompt_version)
    client = genai.Client(api_key=get_gemini_api_key())
    response = client.models.generate_content(
        model=model,
        contents=prompt_config["user_prompt_template"].format(
            output_schema=prompt_config["output_schema"],
            sop_text=sop_text,
        ),
        config=types.GenerateContentConfig(
            system_instruction=prompt_config["system_prompt"],
            temperature=0.4,
            response_mime_type="application/json",
        ),
    )
    payload = safe_json_loads(extract_response_text(response))

    sop_name = str(payload.get("sop_name", "")).strip() or "SOP Training Package"
    training_steps = sanitize_training_steps(payload.get("training_guide") or payload.get("training_steps"))
    summary = remove_summary_overlap(sanitize_summary(payload.get("summary")), training_steps)
    common_mistakes = [
        sentence_case(item) for item in sanitize_string_list(payload.get("common_mistakes", []), limit=3)
    ]
    process_improvements = ensure_process_improvement_coverage(
        sanitize_process_improvements(payload.get("improvements") or payload.get("process_improvements"))
    )
    package = {
        "sop_name": sop_name,
        "summary": summary,
        "training_guide": training_steps,
        "training_steps": training_steps,
        "quiz": sanitize_quiz(payload.get("quiz")),
        "common_mistakes": common_mistakes,
        "improvements": process_improvements,
        "process_improvements": process_improvements,
    }

    package["slides"] = sanitize_slide_deck(payload.get("slides"), f"{sop_name} Deck")
    package["slides"]["slides"] = ensure_slide_slots(package["slides"]["slides"], package)
    if not package["slides"].get("design_notes", {}).get("color_theme"):
        package["slides"]["design_notes"] = build_fallback_slide_deck(sop_name, package)["design_notes"]

    return package


def generate_sop_training_package(
    sop_text: str,
    model: str = DEFAULT_MODEL,
    prompt_version: str | None = None,
) -> Dict[str, Any]:
    if not is_gemini_available():
        raise RuntimeError("Missing GEMINI_API_KEY. Please set your API key.")
    return generate_gemini_training_package(
        sop_text=sop_text,
        model=model,
        prompt_version=prompt_version,
    )


def build_export_text(package: Dict[str, Any]) -> str:
    lines = [package.get("sop_name", "SOP Training Package"), "", "Executive Summary"]

    for bullet in package.get("summary", []):
        lines.append(f"- {bullet}")

    lines.extend(["", "Training Guide"])
    for index, step in enumerate(package.get("training_steps", []), start=1):
        lines.append(f"{index}. {step['title']} | Owner: {step['owner']}")
        lines.append(f"   Action: {step['action']}")
        lines.append(f"   Outcome: {step['outcome']}")
        if step.get("decision_trigger"):
            lines.append(f"   Decision: {step['decision_trigger']}")

    lines.extend(["", "Scenario-Based Quiz"])
    for index, item in enumerate(package.get("quiz", []), start=1):
        lines.append(f"{index}. Question: {item['question']}")
        lines.append(f"   Answer: {item['answer']}")
        lines.append(f"   Explanation: {item['explanation']}")

    lines.extend(["", "Common Mistakes"])
    for item in package.get("common_mistakes", []):
        lines.append(f"- {item}")

    lines.extend(["", "Process Improvements"])
    for item in package.get("process_improvements", []):
        lines.append(f"- {item['theme']}: {item['improvement']} ({item['impact']})")

    return "\n".join(lines)


def build_slides_markdown(package: Dict[str, Any]) -> str:
    slides = package.get("slides", {})
    lines = [f"# {slides.get('deck_title', package.get('sop_name', 'Slide Deck'))}", ""]

    for slide in slides.get("slides", []):
        lines.append(f"## Slide {slide['slide_number']}: {slide['title']}")
        lines.append("")
        for bullet in slide.get("bullets", []):
            lines.append(f"- {bullet}")
        if slide.get("visual_suggestions"):
            lines.append("")
            lines.append("Visual Suggestions:")
            for visual in slide["visual_suggestions"]:
                lines.append(f"- {visual}")
        lines.append("")

    design_notes = slides.get("design_notes", {})
    lines.append("## Design Notes")
    lines.append("")
    lines.append(f"- Color Theme: {design_notes.get('color_theme', '')}")
    lines.append(f"- Font Pairing: {design_notes.get('font_pairing', '')}")
    for suggestion in design_notes.get("layout_suggestions", []):
        lines.append(f"- Layout: {suggestion}")

    return "\n".join(lines).strip()


def build_slides_pptx(package: Dict[str, Any]) -> BytesIO:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as exc:
        raise RuntimeError("PowerPoint export is unavailable right now.") from exc

    presentation = Presentation()
    slides = package.get("slides", {}).get("slides", [])

    for slide_data in slides:
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)

        title = sentence_case(slide_data.get("title", "").strip()) or "Slide"
        bullets = sanitize_string_list(slide_data.get("bullets", []), limit=5)

        slide.shapes.title.text = title
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.size = Pt(26)

        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        for index, bullet in enumerate(bullets[:5]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = sentence_case(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            if hasattr(paragraph, "space_after"):
                paragraph.space_after = Pt(6)

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output
